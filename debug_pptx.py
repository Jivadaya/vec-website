from pptx import Presentation
import os

template_path = 'cert-generator/certificate_template.pptx'
if not os.path.exists(template_path):
    print(f"Error: {template_path} not found")
else:
    prs = Presentation(template_path)
    print(f"Template loaded. Total slides: {len(prs.slides)}")
    
    if len(prs.slides) == 0:
        print("Template has no slides! Checking Master Slides / Layouts...")
        for i, layout in enumerate(prs.slide_layouts):
            placeholders = [f"{p.placeholder_format.idx}: {p.name}" for p in layout.placeholders]
            if placeholders:
                print(f"Layout idx {i} has placeholders: {placeholders}")
    else:
        for i, slide in enumerate(prs.slides):
            print(f"Slide {i} placeholders:")
            for shape in slide.placeholders:
                print(f"  ID: {shape.placeholder_format.idx} | Name: {shape.name} | Text: '{shape.text}'")
