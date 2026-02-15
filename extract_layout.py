from pptx import Presentation

def extract_percentages(pptx_path):
    prs = Presentation(pptx_path)
    w = prs.slide_width
    h = prs.slide_height
    
    print(f"Slide Dimensions: {w}x{h}")
    
    master = prs.slide_masters[0]
    layout = master.slide_layouts[0]
    
    for shape in layout.shapes:
        if not shape.has_text_frame: continue
        
        # Calculate percentage positions
        left_pct = (shape.left / w) * 100
        top_pct = (shape.top / h) * 100
        width_pct = (shape.width / w) * 100
        height_pct = (shape.height / h) * 100
        
        print(f"Shape: {shape.name}")
        print(f"  Left: {left_pct:.2f}%, Top: {top_pct:.2f}%")
        print(f"  Width: {width_pct:.2f}%, Height: {height_pct:.2f}%")
        
        if shape.text_frame.paragraphs:
            p = shape.text_frame.paragraphs[0]
            if p.runs:
                run = p.runs[0]
                # Estimate font size relative to slide height
                # 1pt = 12700 EMUs. Slide height in pts = h / 12700
                slide_h_pt = h / 12700
                font_size_pt = run.font.size.pt if run.font.size else 18
                font_pct = (font_size_pt / slide_h_pt) * 100
                print(f"  Font Ref: {font_size_pt}pt (~{font_pct:.2f}% of slide height)")
        print("-" * 30)

if __name__ == "__main__":
    pptx_path = r"d:\website\vec-website\cert-generator\certificate_template.pptx"
    extract_percentages(pptx_path)
