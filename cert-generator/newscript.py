import os, subprocess, psycopg2, time
from pptx import Presentation
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseUpload

# ======================
# CONFIGURATION
# ======================
if os.name == 'nt':  # Windows
    LIBREOFFICE_PATH = r"C:\Program Files\LibreOffice\program\soffice.exe"
    USE_SHELL = True
else:  # Linux (VPS)
    LIBREOFFICE_PATH = "libreoffice"
    USE_SHELL = False

PPTX_TEMPLATE = "certificate_template.pptx"
OUTPUT_DIR = "generated"
ACADEMIC_YEAR = "2025-26"

# Google Drive Config
CREDENTIALS_FILE = "credentials.json"
TOKEN_FILE = "token.json"
SCOPES = ["https://www.googleapis.com/auth/drive"]
DRIVE_FOLDER_ID = "1Ss2155PsLCthZrAFM4g6TjF_9Rz965sP"

SUPABASE_DB = {
    "host": "aws-1-ap-southeast-1.pooler.supabase.com",
    "dbname": "postgres",
    "user": "postgres.qnvhtrfdwihmyjerujqa",
    "password": "RadhaKrishna@!08", 
    "port": 6543
}

# ======================
# INITIALIZE SERVICES
# ======================
LOCK_FILE = os.path.join(OUTPUT_DIR, "script.lock")

# Simple lock check to prevent cron overlaps
if os.path.exists(LOCK_FILE):
    # Check if the process is actually running (simple age check - 30 mins)
    if (time.time() - os.path.getmtime(LOCK_FILE)) < 1800:
        print("⏳ Script is already running. Skipping this turn.")
        exit()
    else:
        os.remove(LOCK_FILE)

# Create lock
with open(LOCK_FILE, "w") as f:
    f.write(str(os.getpid()))

try:
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    # 1. Drive Service
    creds = None
    if os.path.exists(TOKEN_FILE):
        creds = Credentials.from_authorized_user_file(TOKEN_FILE, SCOPES)

    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(CREDENTIALS_FILE, SCOPES)
            creds = flow.run_local_server(port=0)
        with open(TOKEN_FILE, "w") as token:
            token.write(creds.to_json())
    drive_service = build('drive', 'v3', credentials=creds)

    # 2. Database Connection
    conn = psycopg2.connect(**SUPABASE_DB, sslmode='require')
    cur = conn.cursor()

    # ======================
    # MAIN PROCESS
    # ======================
    cur.execute('SELECT "VEC Exam Code", "Name of Student", "Standard/Class", "School Name", "Result Grades", "District" FROM students WHERE "Certificate Flag" = \'NO\'')
    rows = cur.fetchall()
    
    if not rows:
        print("🎉 No pending certificates to generate.")
    else:
        print(f"🚀 Processing {len(rows)} certificates automatically...")

        for row in rows:
            exam_code, name, standard, school, grade, district = row
            print(f"✍️ Processing: {name} ({exam_code})")

            # --- 1. GENERATE PPTX ---
            try:
                prs = Presentation(PPTX_TEMPLATE)
                slide = prs.slides[0] if len(prs.slides) > 0 else prs.slides.add_slide(prs.slide_layouts[0])

                data_map = {
                    10: str(name).upper(), # ALL CAPS as requested
                    11: str(school),
                    14: ACADEMIC_YEAR,
                    15: str(standard),
                    16: f"Secured: {grade}",
                    17: str(exam_code),
                    18: str(district) # Mapping new District placeholder
                }

                for shape in slide.placeholders:
                    idx = shape.placeholder_format.idx
                    if idx in data_map:
                        # Preserve formatting by modifying the first run
                        if shape.has_text_frame:
                            text_frame = shape.text_frame
                            if text_frame.paragraphs:
                                paragraph = text_frame.paragraphs[0]
                                if paragraph.runs:
                                    paragraph.runs[0].text = str(data_map[idx])
                                    for run_idx in range(len(paragraph.runs) - 1, 0, -1):
                                        p = paragraph._p
                                        p.remove(paragraph.runs[run_idx]._r)
                                else:
                                    paragraph.text = str(data_map[idx])
                            else:
                                shape.text = str(data_map[idx])
                        else:
                            shape.text = str(data_map[idx])

                temp_pptx = os.path.abspath(os.path.join(OUTPUT_DIR, f"{exam_code}.pptx"))
                pdf_path = os.path.abspath(os.path.join(OUTPUT_DIR, f"{exam_code}.pdf"))
                prs.save(temp_pptx)
                
                # --- 2. CONVERT TO PDF ---
                subprocess.run([
                    LIBREOFFICE_PATH, "--headless", "--convert-to", "pdf",
                    "--outdir", os.path.abspath(OUTPUT_DIR), temp_pptx
                ], shell=USE_SHELL, capture_output=True)
                
                # --- 3. AUTO-UPLOAD TO DRIVE ---
                file_metadata = {
                    'name': f"{exam_code}_{name.replace(' ', '_')}.pdf",
                    'parents': [DRIVE_FOLDER_ID]
                }
                
                with open(pdf_path, 'rb') as f:
                    media = MediaIoBaseUpload(f, mimetype='application/pdf', resumable=True)
                    uploaded_file = drive_service.files().create(
                        body=file_metadata,
                        media_body=media,
                        fields='id, webViewLink'
                    ).execute()
                
                file_id = uploaded_file.get('id')
                drive_link = uploaded_file.get('webViewLink')

                # --- 4. FIX PERMISSIONS ---
                drive_service.permissions().create(
                    fileId=file_id,
                    body={'type': 'anyone', 'role': 'reader'}
                ).execute()

                # --- 5. UPDATE SUPABASE ---
                cur.execute("""
                    UPDATE students 
                    SET "Certificate Flag" = 'YES', "Link of Certificate" = %s 
                    WHERE "VEC Exam Code" = %s
                """, (drive_link, exam_code))
                conn.commit()

                # --- 6. CLEANUP ---
                if os.path.exists(temp_pptx): os.remove(temp_pptx)
                if os.path.exists(pdf_path): os.remove(pdf_path)
                
                print(f"✅ Success: {exam_code} is live.")
                time.sleep(1) 

            except Exception as e:
                print(f"⚠️ Failed to process {exam_code}: {e}")
                continue

        print(f"\n🚀 MISSION ACCOMPLISHED! {len(rows)} certificates handled.")

except Exception as e:
    print(f"❌ Critical Error: {e}")
finally:
    # Cleanup lock and connections
    if os.path.exists(LOCK_FILE):
        os.remove(LOCK_FILE)
    try:
        cur.close(); conn.close()
    except:
        pass