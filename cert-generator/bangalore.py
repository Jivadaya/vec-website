import os
import pandas as pd
from pptx import Presentation
from pypdf import PdfReader, PdfWriter
import win32com.client
import gc
import time
from tqdm import tqdm

# ======================
# CONFIG
# ======================
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
EXCEL_FILE = r"D:\website\vec-website\cert-generator\VEC Final Score Sheet_updated Feb 16.xlsx"
PPTX_TEMPLATE = os.path.join(BASE_DIR, "Bangalore.pptx")
OUTPUT_DIR = os.path.join(BASE_DIR, "generated_certificates")

ACADEMIC_YEAR = "2025-26"
BATCH_SIZE = 50  # Balanced for speed vs RAM safety

os.makedirs(OUTPUT_DIR, exist_ok=True)

def clean_filename(text):
    return "".join(c for c in str(text) if c.isalnum() or c == " ").strip().replace(" ", "_")

# Load Data
print("📖 Reading Excel...")
df = pd.read_excel(EXCEL_FILE)
total_students = len(df)

# Start PowerPoint Engine
print("⚡ Launching PowerPoint...")
try:
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = True 
except Exception as e:
    print(f"❌ Error: {e}")
    exit()

pbar = tqdm(total=total_students, desc="Generating", unit="cert")

for batch_start in range(0, total_students, BATCH_SIZE):
    batch_end = min(batch_start + BATCH_SIZE, total_students)
    
    current_batch_meta = []
    
    # Pre-check batch to skip existing files
    for i in range(batch_start, batch_end):
        row = df.iloc[i]
        dist = str(row.get("District", "Unknown")).strip().title()
        sch = str(row.get("School Name", "Unknown")).strip().title()
        nm = str(row.get("Name of Student", "Unknown")).strip().title()
        
        fname = f"VEC_{clean_filename(dist)}_{clean_filename(sch)}_{clean_filename(nm)}.pdf"
        fpath = os.path.join(OUTPUT_DIR, fname)
        
        if os.path.exists(fpath):
            pbar.update(1)
            continue
            
        current_batch_meta.append({
            "row": row, 
            "filename": fname,
            "data": (nm, sch, dist, str(row.get("Standard/Class", "")), str(row.get("Marks", "")), str(row.get("VEC Exam Code", "")))
        })

    if not current_batch_meta:
        continue

    # 1. Build PPTX for this batch
    timestamp = int(time.time())
    temp_ppt = os.path.abspath(os.path.join(OUTPUT_DIR, f"tmp_{timestamp}.pptx"))
    temp_pdf = os.path.abspath(os.path.join(OUTPUT_DIR, f"tmp_{timestamp}.pdf"))

    prs = Presentation(PPTX_TEMPLATE)
    slide_layout = prs.slide_layouts[0]
    
    # Wipe template slides to prevent mismatch
    rId_list = [s.rId for s in prs.slides]
    for rId in rId_list: prs.part.drop_rel(rId)
    prs.slides._sldIdLst.clear()

    for item in current_batch_meta:
        nm, sch, dist, std, marks, code = item["data"]
        slide = prs.slides.add_slide(slide_layout)
        
        data_map = {10: nm, 11: sch, 14: ACADEMIC_YEAR, 15: std, 16: marks, 17: code, 18: dist}
        
        for shape in slide.placeholders:
            idx = shape.placeholder_format.idx
            if idx in data_map:
                shape.text = data_map[idx]

    prs.save(temp_ppt)

    # 2. Export Batch to PDF
    try:
        deck = powerpoint.Presentations.Open(temp_ppt, WithWindow=False)
        deck.SaveAs(temp_pdf, 32)
        deck.Close()
        
        # 3. Split PDF into your specific format
        reader = PdfReader(temp_pdf)
        for idx, page in enumerate(reader.pages):
            writer = PdfWriter()
            writer.add_page(page)
            
            final_path = os.path.join(OUTPUT_DIR, current_batch_meta[idx]["filename"])
            with open(final_path, "wb") as f:
                writer.write(f)
            pbar.update(1)
            
    except Exception as e:
        print(f"\n⚠️ Batch Error: {e}")
    
    # Cleanup
    if os.path.exists(temp_ppt): os.remove(temp_ppt)
    if os.path.exists(temp_pdf): os.remove(temp_pdf)
    gc.collect()

pbar.close()
powerpoint.Quit()
print(f"\n✅ SUCCESS! All individual certificates are in: {OUTPUT_DIR}")